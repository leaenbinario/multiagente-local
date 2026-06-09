import os
import re
import shutil
from pathlib import Path
import email
from email import policy
import hashlib
import logging

import fitz
import chromadb
from chromadb.utils.embedding_functions.ollama_embedding_function import (
    OllamaEmbeddingFunction,
)

from docx import Document
from odf import text, teletype
from odf.opendocument import load
from odf.draw import Page
from odf.text import P as OdfP

from pptx import Presentation

from pdf2image import convert_from_path
from PIL import Image
import pytesseract


logger = logging.getLogger("ingesta_forense")
if not logger.handlers:
    logging.basicConfig(
        level=logging.INFO,
        format="%(message)s"
    )

logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)
logging.getLogger("chromadb").setLevel(logging.WARNING)

BASE_DIR = Path(__file__).resolve().parent

CARPETA_DOCS = BASE_DIR.parent / "shared-data" / "documentos_forenses"
CARPETA_PROCESADOS = BASE_DIR.parent / "shared-data" / "procesados_forense"

COLLECTION_NAME = os.getenv("FORENSE_COLLECTION", "documentos_forense")

CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "200"))

EXT_TXT = {".txt"}
EXT_PDF = {".pdf"}
EXT_EML = {".eml"}
EXT_DOCX = {".docx"}
EXT_ODT = {".odt"}
EXT_ODP = {".odp"}
EXT_DOC = {".doc"}
EXT_PPTX = {".pptx"}

CHROMA_HOST = os.getenv("CHROMA_HOST", "chromadb")
CHROMA_PORT = int(os.getenv("CHROMA_PORT", "8000"))
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text:latest")
OLLAMA_TIMEOUT = int(os.getenv("OLLAMA_TIMEOUT", "300"))
BATCH_SIZE = int(os.getenv("BATCH_SIZE", "10"))


def leer_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def ocr_pdf(path: Path, lang: str = "spa") -> str:
    try:
        pages = convert_from_path(str(path), dpi=300)
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo convertir {path.name} a imágenes para OCR: {e}")
        return ""

    textos_paginas = []
    for i, img in enumerate(pages, start=1):
        try:
            if not isinstance(img, Image.Image):
                img = Image.fromarray(img)
            texto = pytesseract.image_to_string(img, lang=lang)
            if texto.strip():
                textos_paginas.append(f"[PÁGINA {i}]\n{texto}")
        except Exception as e:
            print(f"  [ADVERTENCIA] Error de OCR en {path.name}, página {i}: {e}")

    return "\n\n".join(textos_paginas)


def leer_pdf(path: Path) -> str:
    texto_final = ""
    try:
        textos = []
        with fitz.open(str(path)) as doc:
            for i, page in enumerate(doc, start=1):
                try:
                    texto_pagina = page.get_text() or ""
                    if texto_pagina.strip():
                        textos.append(texto_pagina)
                except Exception as e:
                    print(f"  [ADVERTENCIA] Error extrayendo texto en {path.name}, página {i}: {e}")
        texto_final = "\n".join(textos).strip()
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir/leer PDF {path.name}: {e}")

    if not texto_final:
        print(f"  [INFO] {path.name} parece escaneado o sin texto embebido. Intentando OCR...")
        texto_final = ocr_pdf(path, lang="spa")

    return texto_final


def leer_eml(path: Path) -> str:
    with open(path, "rb") as f:
        msg = email.message_from_bytes(f.read(), policy=policy.default)

    partes = []

    subject = msg.get("subject", "")
    if subject:
        partes.append(f"Asunto: {subject}")

    from_ = msg.get("from", "")
    if from_:
        partes.append(f"De: {from_}")

    to = msg.get("to", "")
    if to:
        partes.append(f"Para: {to}")

    body_texts = []
    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            if ctype in {"text/plain", "text/html"}:
                try:
                    body_texts.append(part.get_content())
                except Exception:
                    pass
    else:
        if msg.get_content_type() in {"text/plain", "text/html"}:
            try:
                body_texts.append(msg.get_content())
            except Exception:
                pass

    if body_texts:
        partes.append("\n".join(body_texts))

    return "\n\n".join(partes)


def leer_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir DOCX {path.name}: {e}")
        return ""

    parrafos = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(parrafos)


def leer_odt(path: Path) -> str:
    try:
        odt_doc = load(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir ODT {path.name}: {e}")
        return ""

    textos = []
    for elem in odt_doc.getElementsByType(text.P):
        t = teletype.extractText(elem)
        if t.strip():
            textos.append(t)

    return "\n".join(textos)


def leer_odp(path: Path) -> str:
    try:
        odp_doc = load(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir ODP {path.name}: {e}")
        return ""

    diapositivas_texto = []

    for i, page in enumerate(odp_doc.getElementsByType(Page), start=1):
        parrafos = []
        for p in page.getElementsByType(OdfP):
            t = teletype.extractText(p)
            if t.strip():
                parrafos.append(t)

        texto_diapo = "\n".join(parrafos).strip()
        if texto_diapo:
            diapositivas_texto.append(f"[DIAPOSITIVA {i}]\n{texto_diapo}")

    return "\n\n".join(diapositivas_texto)


def leer_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir PPTX {path.name}: {e}")
        return ""

    diapositivas_texto = []

    for i, slide in enumerate(prs.slides, start=1):
        textos_cajas = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    texto_parrafo = "".join(run.text for run in paragraph.runs).strip()
                    if texto_parrafo:
                        textos_cajas.append(texto_parrafo)

        texto_diapo = "\n".join(textos_cajas).strip()
        if texto_diapo:
            diapositivas_texto.append(f"[DIAPOSITIVA {i}]\n{texto_diapo}")

    return "\n\n".join(diapositivas_texto)


def extraer_texto(path: Path) -> str:
    ext = path.suffix.lower()

    if ext in EXT_TXT:
        return leer_txt(path)
    if ext in EXT_PDF:
        return leer_pdf(path)
    if ext in EXT_EML:
        return leer_eml(path)
    if ext in EXT_DOCX:
        return leer_docx(path)
    if ext in EXT_ODT:
        return leer_odt(path)
    if ext in EXT_ODP:
        return leer_odp(path)
    if ext in EXT_PPTX:
        return leer_pptx(path)
    if ext in EXT_DOC:
        print(f"  [ADVERTENCIA] {path.name} es un .doc antiguo. Convierte a .docx o .pdf antes de ingestar.")
        return ""

    return ""


def normalizar_texto(texto: str) -> str:
    texto = texto.replace("\r\n", "\n").replace("\r", "\n")
    texto = re.sub(r"[ \t]+", " ", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto.strip()


def infer_title(path: Path, texto: str) -> str:
    lineas = [l.strip() for l in texto.splitlines() if l.strip()]
    candidatas = lineas[:10]

    for linea in candidatas:
        if 8 <= len(linea) <= 220:
            if not linea.lower().startswith(("de:", "para:", "asunto:")):
                return linea

    return path.stem.replace("_", " ").strip()[:220]


def infer_document_type(texto: str, path: Path) -> str:
    t = texto[:5000].lower()
    if re.search(r"\bacordada\b", t):
        return "acordada"
    if re.search(r"\bley\b", t):
        return "ley"
    if re.search(r"\bdecreto\b", t):
        return "decreto"
    if re.search(r"\bresoluci[oó]n\b", t):
        return "resolucion"
    if re.search(r"\bexpte\.?\b|\bexpediente\b|\bcausa\b", t):
        return "expediente"
    if re.search(r"\bpericia\b|\binforme pericial\b", t):
        return "pericia"
    if path.suffix.lower() == ".eml":
        return "email"
    if path.suffix.lower() in {".pptx", ".odp"}:
        return "presentacion"
    return "documento"


def extract_legal_reference(texto: str) -> str:
    patrones = [
        r"\bAcordada\s+\d{1,5}\b",
        r"\bLey\s+\d{1,6}\b",
        r"\bDecreto\s+\d{1,6}(?:/\d{2,4})?\b",
        r"\bResoluci[oó]n\s+\d{1,6}(?:/\d{2,4})?\b",
        r"\bExpte\.?\s*[:\-]?\s*[\w./-]+\b",
        r"\bExpediente\s*[:\-]?\s*[\w./-]+\b",
        r"\bCausa\s*[:\-]?\s*[\w./ -]+\b",
    ]
    for patron in patrones:
        m = re.search(patron, texto, re.IGNORECASE)
        if m:
            return m.group(0).strip()
    return ""


def extract_year(texto: str) -> str:
    matches = re.findall(r"\b(19\d{2}|20\d{2})\b", texto[:5000])
    return matches[0] if matches else ""


def infer_jurisdiction(texto: str) -> str:
    t = texto[:5000].lower()
    if "provincia de buenos aires" in t:
        return "Provincia de Buenos Aires"
    if "ciudad autónoma de buenos aires" in t or "caba" in t:
        return "CABA"
    if "república argentina" in t or "argentina" in t:
        return "Argentina"
    return ""


def infer_issuing_body(texto: str) -> str:
    t = texto[:5000]
    patrones = [
        r"Corte Suprema de Justicia de la Provincia de Buenos Aires",
        r"Suprema Corte de Justicia de la Provincia de Buenos Aires",
        r"Corte Suprema de Justicia de la Naci[oó]n",
        r"Banco Central de la Rep[uú]blica Argentina",
        r"Superintendencia de Bancos e Instituciones Financieras",
    ]
    for patron in patrones:
        m = re.search(patron, t, re.IGNORECASE)
        if m:
            return m.group(0).strip()
    return ""


def build_summary(title: str, document_type: str, legal_reference: str, texto: str) -> str:
    limpio = " ".join(texto.split())
    frases = re.split(r"(?<=[\.\!\?])\s+", limpio)
    frases = [f.strip() for f in frases if len(f.strip()) > 20]
    cuerpo = " ".join(frases[:5])[:1400]

    partes = [f"Título: {title}", f"Tipo: {document_type}"]
    if legal_reference:
        partes.append(f"Referencia: {legal_reference}")
    partes.append(f"Resumen: {cuerpo}")
    return "\n".join(partes).strip()


def chunk_text(texto: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    if overlap >= size:
        raise ValueError("CHUNK_OVERLAP debe ser menor que CHUNK_SIZE")

    texto = normalizar_texto(texto)
    if not texto:
        return []

    bloques = [b.strip() for b in texto.split("\n\n") if b.strip()]
    chunks = []
    actual = ""

    for bloque in bloques:
        if len(bloque) > size:
            if actual.strip():
                chunks.append(actual.strip())
                actual = ""

            start = 0
            while start < len(bloque):
                end = start + size
                parte = bloque[start:end].strip()
                if parte:
                    chunks.append(parte)
                if end >= len(bloque):
                    break
                start = end - overlap
            continue

        candidato = f"{actual}\n\n{bloque}".strip() if actual else bloque
        if len(candidato) <= size:
            actual = candidato
        else:
            if actual.strip():
                chunks.append(actual.strip())
            actual = bloque

    if actual.strip():
        chunks.append(actual.strip())

    return chunks


def get_collection():
    client = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
    embedding_fn = OllamaEmbeddingFunction(
        url=OLLAMA_BASE_URL,
        model_name=EMBED_MODEL,
        timeout=OLLAMA_TIMEOUT,
    )
    return client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=embedding_fn,
    )


def hash_text(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


def existing_hashes_en_batch(collection, hashes):
    existentes = set()

    for h in hashes:
        try:
            res = collection.get(where={"hash": h}, include=["metadatas"])
            metas = res.get("metadatas") or []
            for md in metas:
                if md and md.get("hash") == h:
                    existentes.add(h)
        except Exception as e:
            print(f"  [ADVERTENCIA] No se pudo verificar hash {h}: {e}")

    return existentes


def build_base_metadata(path: Path, texto: str, total_chunks: int):
    title = infer_title(path, texto)
    document_type = infer_document_type(texto, path)
    legal_reference = extract_legal_reference(texto)
    year = extract_year(texto)
    jurisdiction = infer_jurisdiction(texto)
    issuing_body = infer_issuing_body(texto)

    return {
        "source": path.name,
        "title": title[:240],
        "document_type": document_type[:80],
        "legal_reference": legal_reference[:120],
        "year": year[:4],
        "jurisdiction": jurisdiction[:120],
        "issuing_body": issuing_body[:180],
        "doc_type": path.suffix.lower().lstrip("."),
        "total_chunks": total_chunks,
    }


def ingestar_carpeta(carpeta_docs: Path, carpeta_procesados: Path):
    collection = get_collection()
    carpeta_procesados.mkdir(parents=True, exist_ok=True)

    archivos = sorted([p for p in carpeta_docs.iterdir() if p.is_file()])
    if not archivos:
        print(f"No hay archivos en {carpeta_docs}")
        return

    for path in archivos:
        try:
            ext = path.suffix.lower()
            if ext not in (EXT_TXT | EXT_PDF | EXT_EML | EXT_DOCX | EXT_ODT | EXT_ODP | EXT_PPTX | EXT_DOC):
                print(f"Salto {path.name} (extensión no soportada por ahora)")
                continue

            print(f"\nProcesando {path.name} ...")
            texto = extraer_texto(path).strip()
            texto = normalizar_texto(texto)

            if not texto:
                print(f"  No se pudo extraer texto de {path.name}, lo salto.")
                continue

            chunks = chunk_text(texto)
            total_chunks = len(chunks)
            meta_base = build_base_metadata(path, texto, total_chunks)

            summary = build_summary(
                title=meta_base["title"],
                document_type=meta_base["document_type"],
                legal_reference=meta_base["legal_reference"],
                texto=texto,
            )

            print(f"  Título: {meta_base['title']}")
            print(f"  Tipo: {meta_base['document_type']}")
            if meta_base["legal_reference"]:
                print(f"  Referencia: {meta_base['legal_reference']}")
            print(f"  Generados {total_chunks} chunks.")

            insertados_total = 0
            duplicados_total = 0

            summary_hash = hash_text(summary)
            summary_meta = {
                **meta_base,
                "chunk_index": -1,
                "order_in_doc": 0,
                "is_summary": True,
                "hash": summary_hash,
            }

            if summary_hash not in existing_hashes_en_batch(collection, [summary_hash]):
                try:
                    collection.upsert(
                        documents=[summary],
                        ids=[f"{path.name}_summary"],
                        metadatas=[summary_meta],
                    )
                    insertados_total += 1
                    print("  Summary insertado.")
                except Exception as e:
                    print(f"  [ERROR] Falló inserción de summary para {path.name}: {e}")
            else:
                duplicados_total += 1
                print("  Summary ya existía.")

            for start in range(0, total_chunks, BATCH_SIZE):
                end = min(start + BATCH_SIZE, total_chunks)
                batch_chunks = chunks[start:end]

                batch_ids = []
                batch_docs = []
                batch_metadatas = []
                batch_hashes = []

                for i, chunk in enumerate(batch_chunks, start=start):
                    contextualizado = (
                        f"Título: {meta_base['title']}\n"
                        f"Tipo: {meta_base['document_type']}\n"
                        f"Referencia: {meta_base['legal_reference']}\n"
                        f"Jurisdicción: {meta_base['jurisdiction']}\n"
                        f"Emisor: {meta_base['issuing_body']}\n"
                        f"Contenido:\n{chunk}"
                    ).strip()

                    h = hash_text(contextualizado)
                    batch_hashes.append(h)
                    batch_ids.append(f"{path.name}_chunk_{i}")
                    batch_docs.append(contextualizado)
                    batch_metadatas.append(
                        {
                            **meta_base,
                            "chunk_index": i,
                            "order_in_doc": i + 1,
                            "is_summary": False,
                            "hash": h,
                        }
                    )

                existing_hashes = existing_hashes_en_batch(collection, batch_hashes)

                filtered_docs = []
                filtered_ids = []
                filtered_metas = []

                for doc, _id, meta in zip(batch_docs, batch_ids, batch_metadatas):
                    if meta["hash"] in existing_hashes:
                        duplicados_total += 1
                        continue
                    filtered_docs.append(doc)
                    filtered_ids.append(_id)
                    filtered_metas.append(meta)

                if not filtered_docs:
                    print(f"  Batch {start}-{end - 1}: ya existía completo.")
                    continue

                try:
                    collection.upsert(
                        documents=filtered_docs,
                        ids=filtered_ids,
                        metadatas=filtered_metas,
                    )
                    insertados_total += len(filtered_docs)
                    print(f"  Batch {start}-{end - 1}: insertados={len(filtered_docs)}")
                except Exception as e:
                    print(f"  [ERROR] Falló batch {start}-{end - 1} de {path.name}: {e}")
                    continue

            print(
                f"  Resumen: totales={total_chunks + 1} | "
                f"insertados={insertados_total} | duplicados={duplicados_total}"
            )

            destino = carpeta_procesados / path.name
            shutil.move(str(path), destino)
            print(f"  Movido a procesados: {destino.name}")

        except Exception as e:
            print(f"  [ERROR] Falló el archivo {path.name}: {e}")
            continue

    print("\nIngesta completa.")


if __name__ == "__main__":
    print(f"Carpeta de documentos FORENSE: {CARPETA_DOCS.resolve()}")
    print(f"Carpeta de procesados FORENSE: {CARPETA_PROCESADOS.resolve()}")
    print(f"Chroma host: {CHROMA_HOST}:{CHROMA_PORT}")
    print(f"Ollama base URL: {OLLAMA_BASE_URL}")
    print(f"Colección Chroma: {COLLECTION_NAME}")

    if not CARPETA_DOCS.exists():
        print("La carpeta documentos_forenses no existe. Créala y coloca allí .txt, .pdf, .eml, .docx, .odt, .odp, .pptx.")
    else:
        ingestar_carpeta(CARPETA_DOCS, CARPETA_PROCESADOS)