import shutil
from pathlib import Path
import os
import chromadb
import email
from email import policy
import hashlib
from pptx import Presentation

import fitz  # PyMuPDF
from chromadb.utils.embedding_functions.ollama_embedding_function import (
    OllamaEmbeddingFunction,
)

# OCR
from pdf2image import convert_from_path
from PIL import Image
import pytesseract

from docx import Document          # para .docx
from odf import text, teletype     # para .odt
from odf.opendocument import load  # para .odt y .odp
from odf.draw import Page          # para diapositivas .odp
from odf.text import P as OdfP     # párrafos en las diapositivas


# ==============================
# 1) Configuración
# ==============================

# Dentro del contenedor, shared-data está montado en /shared-data
SHARED_DIR = Path(__file__).resolve().parent  # /app

CARPETA_DOCS = Path("/shared-data/documentos_operativa")       # entrada
CARPETA_PROCESADOS = Path("/shared-data/procesados_asistente")

# Colección específica para ASISTENTE
COLLECTION_NAME = "documentos_operativa"

# Tamaño de chunk (caracteres) y solapamiento
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

EXT_TXT = {".txt"}
EXT_PDF = {".pdf"}
EXT_EML = {".eml"}
EXT_DOCX = {".docx"}
EXT_ODT = {".odt"}
EXT_ODP = {".odp"}
EXT_DOC = {".doc"}   # solo para advertir
EXT_PPTX = {".pptx"}

# Config para Chroma y Ollama (igual que en forense/contador)
CHROMA_HOST = os.getenv("CHROMA_HOST", "chromadb")
CHROMA_PORT = int(os.getenv("CHROMA_PORT", "8000"))
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://ollama:11434")
EMBED_MODEL = "nomic-embed-text:latest"


# ==============================
# 2) Utilidades de extracción de texto
# ==============================

def leer_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def ocr_pdf(path: Path, lang: str = "spa") -> str:
    """
    Aplica OCR página por página a un PDF escaneado usando pdf2image + pytesseract.
    """
    try:
        pages = convert_from_path(str(path), dpi=300)
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo convertir {path.name} a imágenes para OCR: {e}")
        return ""

    textos_paginas = []
    for i, img in enumerate(pages):
        try:
            if not isinstance(img, Image.Image):
                img = Image.fromarray(img)
            texto = pytesseract.image_to_string(img, lang=lang)
            if texto.strip():
                textos_paginas.append(f"[PÁGINA {i+1}]\n{texto}")
        except Exception as e:
            print(f"  [ADVERTENCIA] Error de OCR en {path.name}, página {i}: {e}")
            continue

    return "\n\n".join(textos_paginas)


def leer_pdf(path: Path) -> str:
    """
    Lectura de PDF:
      1) Intenta extraer texto con PyMuPDF.
      2) Si no obtiene nada útil, intenta OCR con pytesseract.
    """
    try:
        textos = []
        with fitz.open(str(path)) as doc:
            for i, page in enumerate(doc):
                try:
                    texto_pagina = page.get_text() or ""
                    textos.append(texto_pagina)
                except Exception as e:
                    print(f"  [ADVERTENCIA] Error extrayendo texto en {path.name}, página {i}: {e}")
                    continue
        texto_final = "\n".join(textos).strip()
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir/leer PDF {path.name}: {e}")
        texto_final = ""

    if not texto_final:
        print(f"  [INFO] {path.name} parece escaneado o sin texto embebido. Intentando OCR...")
        texto_final = ocr_pdf(path, lang="spa")

    return texto_final


def leer_eml(path: Path) -> str:
    with open(path, "rb") as f:
        msg = email.message_from_bytes(f.read(), policy=policy.default)

    partes = []

    subject = msg.get("subject", "")
    if subject:
        partes.append(f"Asunto: {subject}")

    from_ = msg.get("from", "")
    if from_:
        partes.append(f"De: {from_}")

    to = msg.get("to", "")
    if to:
        partes.append(f"Para: {to}")

    body_texts = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body_texts.append(part.get_content())
    else:
        if msg.get_content_type() == "text/plain":
            body_texts.append(msg.get_content())

    if body_texts:
        partes.append("\n".join(body_texts))

    return "\n\n".join(partes)


def leer_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir DOCX {path.name}: {e}")
        return ""

    parrafos = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(parrafos)


def leer_odt(path: Path) -> str:
    try:
        odt_doc = load(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir ODT {path.name}: {e}")
        return ""

    textos = []
    for elem in odt_doc.getElementsByType(text.P):
        textos.append(teletype.extractText(elem))
    return "\n".join(textos)


def leer_odp(path: Path) -> str:
    """
    Extrae texto de un .odp usando odfpy: recorre diapositivas y párrafos.
    """
    try:
        odp_doc = load(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir ODP {path.name}: {e}")
        return ""

    diapositivas_texto = []

    for i, page in enumerate(odp_doc.getElementsByType(Page)):
        parrafos = []
        for p in page.getElementsByType(OdfP):
            parrafos.append(teletype.extractText(p))
        texto_diapo = "\n".join([t for t in parrafos if t])
        if texto_diapo.strip():
            diapositivas_texto.append(f"[DIAPOSITIVA {i+1}]\n{texto_diapo}")

    return "\n\n".join(diapositivas_texto)


def leer_pptx(path: Path) -> str:
    """
    Extrae texto de un .pptx usando python-pptx: recorre diapositivas y cuadros de texto.
    """
    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"  [ADVERTENCIA] No se pudo abrir PPTX {path.name}: {e}")
        return ""

    diapositivas_texto = []

    for i, slide in enumerate(prs.slides, start=1):
        textos_cajas = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    texto_parrafo = "".join(run.text for run in paragraph.runs).strip()
                    if texto_parrafo:
                        textos_cajas.append(texto_parrafo)

        texto_diapo = "\n".join(textos_cajas).strip()
        if texto_diapo:
            diapositivas_texto.append(f"[DIAPOSITIVA {i}]\n{texto_diapo}")

    return "\n\n".join(diapositivas_texto)


def extraer_texto(path: Path) -> str:
    ext = path.suffix.lower()

    if ext in EXT_TXT:
        return leer_txt(path)
    if ext in EXT_PDF:
        return leer_pdf(path)
    if ext in EXT_EML:
        return leer_eml(path)
    if ext in EXT_DOCX:
        return leer_docx(path)
    if ext in EXT_ODT:
        return leer_odt(path)
    if ext in EXT_ODP:
        return leer_odp(path)
    if ext in EXT_PPTX:
        return leer_pptx(path)
    if ext in EXT_DOC:
        print(
            f"  [ADVERTENCIA] {path.name} es un .doc antiguo. "
            "Convierte a .docx o .pdf antes de ingestar."
        )
        return ""

    return ""


# ==============================
# 3) Chunking sencillo
# ==============================

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    chunks = []
    start = 0
    length = len(text)

    while start < length:
        end = start + size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap

    return chunks


# ==============================
# 4) Inicializar Chroma + embeddings (ASISTENTE)
# ==============================

def get_collection():
    client = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
    embedding_fn = OllamaEmbeddingFunction(
        url=OLLAMA_BASE_URL,
        model_name=EMBED_MODEL,
    )
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=embedding_fn,
    )
    return collection


# ==============================
# 5) Hash para deduplicación global
# ==============================

def hash_text(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


# ==============================
# 6) Ingesta con deduplicación
# ==============================

def ingestar_carpeta(carpeta_docs: Path, carpeta_procesados: Path):
    collection = get_collection()

    carpeta_procesados.mkdir(parents=True, exist_ok=True)

    archivos = sorted([p for p in carpeta_docs.iterdir() if p.is_file()])
    if not archivos:
        print(f"No hay archivos en {carpeta_docs}")
        return

    for path in archivos:
        ext = path.suffix.lower()
        if (
            ext not in EXT_TXT
            and ext not in EXT_PDF
            and ext not in EXT_EML
            and ext not in EXT_DOCX
            and ext not in EXT_ODT
            and ext not in EXT_ODP
            and ext not in EXT_PPTX
            and ext not in EXT_DOC
        ):
            print(f"Salto {path.name} (extensión no soportada por ahora)")
            continue

        print(f"\nProcesando {path.name} ...")
        texto = extraer_texto(path).strip()
        if not texto:
            print(f"  No se pudo extraer texto de {path.name}, lo salto.")
            continue

        chunks = chunk_text(texto)
        total_chunks = len(chunks)
        print(f"  Generados {total_chunks} chunks.")

        batch_size = 10
        for start in range(0, total_chunks, batch_size):
            end = min(start + batch_size, total_chunks)
            batch_chunks = chunks[start:end]

        batch_ids = []
        batch_metadatas = []
        batch_hashes = []

        for i, chunk in enumerate(batch_chunks, start=start):
            h = hash_text(chunk)
            batch_hashes.append(h)
            batch_ids.append(f"{path.name}_chunk_{i}")
            batch_metadatas.append(
                {"source": path.name, "chunk_index": i, "hash": h}
            )

        existing_hashes = set()
        if batch_hashes:
            existing = collection.get(
                where={"hash": {"$in": batch_hashes}},
                include=["metadatas"],
            )
            if existing and existing.get("metadatas"):
                for md in existing["metadatas"]:
                    if md and "hash" in md:
                        existing_hashes.add(md["hash"])

        filtered_docs = []
        filtered_ids = []
        filtered_metas = []

        for doc, _id, meta in zip(batch_chunks, batch_ids, batch_metadatas):
            if meta["hash"] in existing_hashes:
                continue
            filtered_docs.append(doc)
            filtered_ids.append(_id)
            filtered_metas.append(meta)

        if not filtered_docs:
            print("  Todos los chunks de este batch ya existían. Nada que insertar.")
            continue

        collection.upsert(
            documents=filtered_docs,
            ids=filtered_ids,
            metadatas=filtered_metas,
        )

        print(f"  Ingesta terminada para {path.name} en colección {COLLECTION_NAME}.")

        destino = carpeta_procesados / path.name
        shutil.move(str(path), destino)
        print(f"  Movido a procesados: {destino.name}")

    print("\nIngesta completa.")


# ==============================
# 7) Main
# ==============================

if __name__ == "__main__":
    print(f"Carpeta de documentos:  {CARPETA_DOCS.resolve()}")
    print(f"Carpeta de procesados:  {CARPETA_PROCESADOS.resolve()}")

    if not CARPETA_DOCS.exists():
        print("La carpeta documentos_operativa no existe. Créala y coloca allí .txt, .pdf, .eml, .docx, .odt, .odp, .pptx.")
    else:
        ingestar_carpeta(CARPETA_DOCS, CARPETA_PROCESADOS)